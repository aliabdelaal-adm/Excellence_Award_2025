#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
سكريبت لاختبار والتحقق من ملفات PowerPoint
"""

import os
from pptx import Presentation

def test_powerpoint_file(filepath):
    """اختبار ملف PowerPoint واحد"""
    try:
        prs = Presentation(filepath)
        
        print(f"\n{'='*80}")
        print(f"📄 ملف: {os.path.basename(filepath)}")
        print(f"{'='*80}")
        
        print(f"📊 عدد الشرائح: {len(prs.slides)}")
        print(f"📐 حجم الشريحة: {prs.slide_width/914400:.2f} × {prs.slide_height/914400:.2f} بوصة")
        
        print("\n🎯 محتوى الشرائح:")
        for i, slide in enumerate(prs.slides, 1):
            print(f"\nالشريحة {i}:")
            
            # عد عناصر الشريحة
            text_boxes = 0
            total_text = 0
            
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text_boxes += 1
                    text = shape.text_frame.text.strip()
                    if text:
                        total_text += len(text)
                        # عرض أول 60 حرف من النص
                        preview = text[:60] + "..." if len(text) > 60 else text
                        print(f"  ✓ نص: {preview}")
            
            print(f"  📝 إجمالي: {text_boxes} صندوق نص، {total_text} حرف")
        
        print(f"\n✅ الملف صحيح وجاهز للاستخدام")
        return True
        
    except Exception as e:
        print(f"\n❌ خطأ في الملف {filepath}: {str(e)}")
        return False

def main():
    """الدالة الرئيسية"""
    print("🔍 بدء اختبار ملفات PowerPoint...")
    
    ppt_dir = "PowerPoint_Presentations"
    
    if not os.path.exists(ppt_dir):
        print(f"❌ المجلد {ppt_dir} غير موجود!")
        return
    
    ppt_files = [f for f in os.listdir(ppt_dir) if f.endswith('.pptx')]
    
    if not ppt_files:
        print(f"❌ لا توجد ملفات PowerPoint في المجلد {ppt_dir}!")
        return
    
    print(f"📁 تم العثور على {len(ppt_files)} ملف PowerPoint")
    
    success_count = 0
    fail_count = 0
    
    for ppt_file in sorted(ppt_files):
        filepath = os.path.join(ppt_dir, ppt_file)
        if test_powerpoint_file(filepath):
            success_count += 1
        else:
            fail_count += 1
    
    print(f"\n{'='*80}")
    print(f"📈 نتائج الاختبار:")
    print(f"✅ ملفات ناجحة: {success_count}")
    print(f"❌ ملفات فاشلة: {fail_count}")
    print(f"📊 إجمالي: {success_count + fail_count}")
    print(f"{'='*80}")

if __name__ == "__main__":
    main()
