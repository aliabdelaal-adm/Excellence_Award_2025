#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
مولد عروض PowerPoint للمشاريع
يقوم بإنشاء ملف PowerPoint احترافي لكل مشروع/إنجاز بناءً على بيانات plan-data.json
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

# ألوان بلدية أبوظبي
ADM_BLUE = RGBColor(0, 102, 204)  # #0066cc
ADM_GREEN = RGBColor(0, 166, 81)  # #00a651
ADM_GOLD = RGBColor(255, 215, 0)  # #ffd700
DARK_GRAY = RGBColor(44, 62, 80)  # #2c3e50
LIGHT_GRAY = RGBColor(127, 140, 141)  # #7f8c8d

def load_project_data():
    """تحميل بيانات المشاريع من plan-data.json"""
    with open('plan-data.json', 'r', encoding='utf-8') as f:
        return json.load(f)

def create_title_slide(prs, project):
    """إنشاء صفحة العنوان"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # خلفية زرقاء
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ADM_BLUE
    
    # إضافة الأيقونة والعنوان
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p = title_frame.paragraphs[0]
    p.text = f"{project.get('icon', '🎯')} {project['title']}"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # الوصف
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    
    p = desc_frame.paragraphs[0]
    p.text = project.get('description', '')
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # معلومات المؤلف والتاريخ
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
    footer_frame = footer_box.text_frame
    
    p = footer_frame.paragraphs[0]
    p.text = "د. علي عبدالعال | إدارة الرفق بالحيوان | بلدية مدينة أبوظبي"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # التاريخ
    date_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    
    p = date_frame.paragraphs[0]
    p.text = datetime.now().strftime("%Y/%m/%d")
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_section_slide(prs, title, content_list, icon="📊"):
    """إضافة شريحة قسم مع نقاط"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # العنوان
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    
    p = title_frame.paragraphs[0]
    p.text = f"{icon} {title}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ADM_BLUE
    p.alignment = PP_ALIGN.RIGHT
    
    # خط فاصل
    line = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(1.4), Inches(9), Inches(0)
    )
    line.line.color.rgb = ADM_GREEN
    line.line.width = Pt(3)
    
    # المحتوى
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for item in content_list:
        p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(6)
        p.alignment = PP_ALIGN.RIGHT

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """إضافة شريحة بعمودين"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # العنوان
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ADM_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # العمود الأيمن
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ADM_GREEN
    p.alignment = PP_ALIGN.RIGHT
    
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.alignment = PP_ALIGN.RIGHT
    
    # العمود الأيسر
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ADM_GREEN
    p.alignment = PP_ALIGN.RIGHT
    
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.alignment = PP_ALIGN.RIGHT

def add_closing_slide(prs):
    """إضافة شريحة ختامية"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # خلفية خضراء
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ADM_GREEN
    
    # رسالة الشكر
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    thanks_frame = thanks_box.text_frame
    
    p = thanks_frame.paragraphs[0]
    p.text = "شكراً لاهتمامكم"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # معلومات التواصل
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
    contact_frame = contact_box.text_frame
    
    contact_info = [
        "د. علي عبدالعال - طبيب بيطري",
        "إدارة الرفق بالحيوان",
        "بلدية مدينة أبوظبي",
        "📧 Ali.Abdelaal@adm.gov.ae",
        "📱 0581187777"
    ]
    
    for i, info in enumerate(contact_info):
        p = contact_frame.add_paragraph() if i > 0 else contact_frame.paragraphs[0]
        p.text = info
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

def create_project_presentation(project, data):
    """إنشاء عرض تقديمي كامل للمشروع"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. صفحة العنوان
    create_title_slide(prs, project)
    
    # 2. نظرة عامة وأهداف المشروع
    overview_content = [
        f"الفئة: {project.get('category', 'مشروع')}",
        f"الهدف: {project.get('description', '')}",
    ]
    if 'goals' in project:
        overview_content.append(f"المهام: {project['goals']}")
    
    add_section_slide(prs, "نظرة عامة وأهداف المشروع", overview_content, "🎯")
    
    # 3. التقنيات المستخدمة
    if 'technologies' in project and project['technologies']:
        tech_list = project['technologies']
        add_section_slide(prs, "التقنيات والأدوات المستخدمة", tech_list, "💻")
    
    # 4. الفوائد والقيمة المضافة
    if 'benefits' in project and project['benefits']:
        benefits_list = project['benefits']
        # تقسيم الفوائد إلى شرائح إذا كانت كثيرة
        chunk_size = 6
        for i in range(0, len(benefits_list), chunk_size):
            chunk = benefits_list[i:i+chunk_size]
            title = "الفوائد والقيمة المضافة" if i == 0 else f"الفوائد والقيمة المضافة (تابع {i//chunk_size + 1})"
            add_section_slide(prs, title, chunk, "💎")
    
    # 5. الربط مع إدارة الرفق بالحيوان
    animal_welfare_context = [
        "يساهم في تحقيق رؤية بلدية أبوظبي للتميز المؤسسي",
        "يعزز كفاءة العمل في إدارة الرفق بالحيوان",
        "يدعم الأهداف الاستراتيجية للإدارة",
        "يوفر حلول تقنية متقدمة للعمليات اليومية",
        "يحسن جودة الخدمات المقدمة للمجتمع",
        "يساهم في رفع مستوى الأداء المؤسسي"
    ]
    add_section_slide(prs, "الربط مع إدارة الرفق بالحيوان", animal_welfare_context, "🐾")
    
    # 6. الأهداف التشغيلية والاستراتيجية
    operational_goals = [
        "تحسين الكفاءة التشغيلية للإدارة",
        "تقليل الوقت والجهد المبذول في العمليات",
        "رفع مستوى الدقة وتقليل الأخطاء",
        "تعزيز التواصل والتعاون بين الفرق"
    ]
    
    strategic_goals = [
        "تحقيق التحول الرقمي الشامل",
        "تطبيق معايير التميز العالمية",
        "دعم رؤية أبوظبي 2030",
        "تعزيز مكانة بلدية أبوظبي الريادية"
    ]
    
    add_two_column_slide(prs, "الأهداف التشغيلية والاستراتيجية", 
                         "الأهداف الاستراتيجية 🎯", strategic_goals,
                         "الأهداف التشغيلية ⚙️", operational_goals)
    
    # 7. البيانات الداعمة والإحصائيات
    if 'benefits' in project and project['benefits']:
        # استخراج الأرقام والإحصائيات من الفوائد
        stats = []
        for benefit in project['benefits']:
            if any(char.isdigit() for char in benefit):
                stats.append(benefit)
        
        if stats:
            add_section_slide(prs, "البيانات الداعمة والإحصائيات", stats, "📊")
    
    # 8. الأدلة والوثائق المساندة
    evidence = [
        f"ملف المشروع: {project.get('link', 'غير متوفر')}",
        "التوثيق الكامل متوفر في المنصة الإلكترونية",
        "عروض تفاعلية HTML متاحة للمراجعة",
        "بيانات تفصيلية في plan-data.json",
        "تقارير تحليلية في المجلدات التفاعلية",
        "معايير دولية موثقة في best-practices-summary.md"
    ]
    add_section_slide(prs, "الأدلة والوثائق المساندة", evidence, "📋")
    
    # 9. التأثير والنتائج المتوقعة
    impact = [
        "تحسين مستوى الخدمات المقدمة",
        "زيادة رضا المستفيدين الداخليين والخارجيين",
        "توفير الموارد والوقت والجهد",
        "تعزيز صورة البلدية كمؤسسة رائدة",
        "المساهمة في تحقيق أهداف التنمية المستدامة",
        "دعم التحول الرقمي الوطني"
    ]
    add_section_slide(prs, "التأثير والنتائج المتوقعة", impact, "🎯")
    
    # 10. صفحة ختامية
    add_closing_slide(prs)
    
    return prs

def main():
    """الدالة الرئيسية"""
    print("🚀 بدء إنشاء ملفات PowerPoint للمشاريع...")
    
    # تحميل البيانات
    data = load_project_data()
    
    # إنشاء مجلد للعروض التقديمية
    output_dir = "PowerPoint_Presentations"
    os.makedirs(output_dir, exist_ok=True)
    
    # جمع كل المشاريع والإنجازات
    all_items = []
    
    # المشاريع الرئيسية
    if 'projects' in data:
        all_items.extend(data['projects'])
    
    # ملف الإنجازات
    if 'portfolio' in data:
        all_items.extend(data['portfolio'])
    
    # إنشاء عرض تقديمي لكل مشروع
    created_count = 0
    for item in all_items:
        try:
            # تخطي العناصر بدون عنوان أو مع روابط فارغة
            if not item.get('title') or item.get('link') == '#':
                continue
            
            print(f"\n📄 إنشاء عرض تقديمي: {item['title']}")
            
            prs = create_project_presentation(item, data)
            
            # حفظ الملف
            safe_title = item['title'].replace('/', '-').replace('\\', '-')
            filename = f"{output_dir}/{safe_title}.pptx"
            prs.save(filename)
            
            print(f"✅ تم حفظ الملف: {filename}")
            created_count += 1
            
        except Exception as e:
            print(f"❌ خطأ في إنشاء عرض تقديمي لـ {item.get('title', 'Unknown')}: {str(e)}")
    
    print(f"\n✨ تم الانتهاء! تم إنشاء {created_count} ملف PowerPoint في المجلد {output_dir}")
    print(f"📂 المسار الكامل: {os.path.abspath(output_dir)}")

if __name__ == "__main__":
    main()
